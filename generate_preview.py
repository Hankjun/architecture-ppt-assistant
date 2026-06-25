#!/usr/bin/env python3
"""
建筑汇报 PPT 自动助手 — HTML 多页预览生成器
根据结构化数据生成可预览的 HTML 演示页面，支持多种页面模板。

用法:
  python3 generate_preview.py                          # 默认示例
  python3 generate_preview.py --json project.json       # 自定义数据
  python3 generate_preview.py --theme light             # 浅色主题
  python3 generate_preview.py --pptx output.pptx        # 导出 PPTX
"""

import json, argparse, os, subprocess, sys
from datetime import datetime

# =============================================================================
# 页面模板
# =============================================================================

PAGE_TYPES = {
    "cover":       "封面 — 标题、副标题、元信息",
    "toc":         "目录 — 汇报框架列表",
    "content":     "内容页 — 标题+正文+数据统计(4格)",
    "split":       "双栏 — 左右两栏策略/对比",
    "comparison":  "对比 — 新旧/AB对比，带标签",
    "gallery":     "图库 — 图片网格展示",
    "quote":       "引用 — 视觉突出的引言页",
    "end":         "结束页 — 感谢",
}

DEFAULT_SECTIONS = {
    "title": "某新城核心区城市设计",
    "subtitle": "概念规划方案汇报",
    "client": "某市自然资源和规划局",
    "date": datetime.now().strftime("%Y年%m月"),
    "theme": "dark",
    "pages": [
        {"type": "cover", "title": "某新城核心区城市设计", "subtitle": "概念规划方案汇报",
         "meta": "某市自然资源和规划局 · 2025"},
        {"type": "toc", "title": "汇报框架", "items": [
            "01 / 项目区位", "02 / 现状分析", "03 / 设计策略",
            "04 / 总图方案", "05 / 建筑分析", "06 / 效果展示"]},
        {"type": "content", "title": "01 / 项目区位",
         "subtitle": "基地位于新城核心区，交通枢纽辐射范围 3km",
         "body": "项目位于新城核心区东侧，紧邻城市主干道及地铁3号线站点，周边配套完善。基地总面积约12.6公顷，现状以空地及低效工业用地为主。",
         "stats": [{"label": "基地面积", "value": "12.6 ha"}, {"label": "容积率", "value": "2.8"},
                   {"label": "限高", "value": "80 m"}, {"label": "绿地率", "value": "≥ 35%"}]},
        {"type": "comparison", "title": "方案对比",
         "subtitle": "两版总图方案比较",
         "left_label": "方案 A · 集中布局",
         "left_body": "核心功能集中布置在地铁站点周边，形成高密度 TOD 核心区，商业与办公垂直混合。",
         "right_label": "方案 B · 组团分散",
         "right_body": "功能分散为三个独立组团，各组团围绕小型开放空间布局，通过慢行系统串联。"},
        {"type": "quote", "quote": "设计不是让事物看起来怎样，\n而是让事物如何工作。",
         "attribution": "——《建筑的永恒之道》"},
        {"type": "content", "title": "04 / 总图方案",
         "subtitle": "一轴两带三组团 · 功能复合的城市核心",
         "body": "总图以\"一轴两带三组团\"为结构：南北向城市活力轴串联各功能组团；东西两条生态绿带渗透基地；形成 TOD 综合组团、品质居住组团、滨水商业组团三大板块。",
         "stats": [{"label": "总建筑面积", "value": "35.2 万㎡"}, {"label": "商业", "value": "8.6 万㎡"},
                   {"label": "住宅", "value": "18.4 万㎡"}, {"label": "公建配套", "value": "8.2 万㎡"}]},
        {"type": "gallery", "title": "效果展示", "subtitle": "方案效果图预览",
         "images": [
             {"caption": "鸟瞰效果图", "color": "#c96442"},
             {"caption": "主入口透视图", "color": "#4d4c48"},
             {"caption": "滨水商业街", "color": "#30302e"},
             {"caption": "核心区广场", "color": "#1a1a19"}]},
        {"type": "end", "title": "感谢聆听", "subtitle": "敬请指正 · 期待合作"}
    ]
}


# =============================================================================
# HTML 生成
# =============================================================================

def generate_html(data):
    pages_html = "".join(_render_page(p, i, data.get("theme", "dark")) for i, p in enumerate(data["pages"]))
    theme = data.get("theme", "dark")

    if theme == "light":
        bg, text, sec_bg, sec_text, card_bg, border, accent = "#faf9f5", "#141413", "#f0eee6", "#4d4c48", "#ffffff", "#e8e6dc", "#c96442"
        subtitle_color, stat_bg = "#87867f", "#f5f4ed"
    else:
        bg, text, sec_bg, sec_text, card_bg, border, accent = "#141413", "#faf9f5", "#1a1a19", "#87867f", "#1a1a19", "#30302e", "#c96442"
        subtitle_color, stat_bg = "#87867f", "#1a1a19"

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{data['title']} · 预览</title>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
  <style>
    *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
    html {{ scroll-behavior: smooth; scroll-snap-type: y proximity; }}
    body {{
      font-family: 'Inter', system-ui, -apple-system, sans-serif;
      background: {bg}; color: {text};
      -webkit-font-smoothing: antialiased;
    }}
    .slide {{
      min-height: 100vh; padding: 4rem 3rem;
      display: flex; flex-direction: column; justify-content: center;
      scroll-snap-align: start;
      border-bottom: 1px solid {border};
      max-width: 960px; margin: 0 auto;
    }}
    .slide:last-child {{ border-bottom: none; }}
    .slide-label {{
      font-size: 0.7rem; font-weight: 500; letter-spacing: 0.12px;
      text-transform: uppercase; color: #87867f; margin-bottom: 0.5rem;
    }}
    .slide h1 {{
      font-family: Georgia, 'Times New Roman', Times, serif;
      font-size: clamp(2rem,4vw,3.5rem); font-weight: 500;
      line-height: 1.15; color: {text}; margin-bottom: 0.75rem;
    }}
    .slide h2 {{
      font-family: Georgia, 'Times New Roman', Times, serif;
      font-size: clamp(1.4rem,2.5vw,2rem); font-weight: 500;
      line-height: 1.25; color: {text}; margin-bottom: 0.75rem;
    }}
    .slide h3 {{
      font-family: Georgia, 'Times New Roman', Times, serif;
      font-size: 1.15rem; font-weight: 500;
      color: {accent}; margin-bottom: 0.5rem;
    }}
    .slide p {{ color: {subtitle_color}; line-height: 1.8; font-size: 0.95rem; max-width: 640px; }}
    .subtitle {{ color: {subtitle_color}; font-size: 1rem; line-height: 1.6; margin-bottom: 1.5rem; max-width: 600px; }}
    .cover {{ text-align: center; }}
    .cover .meta {{ color: #87867f; font-size: 0.85rem; margin-top: 1.5rem; }}
    .accent-line {{ width: 40px; height: 3px; background: {accent}; margin: 1.5rem auto; border-radius: 2px; }}
    .toc-list {{ list-style: none; display: flex; flex-direction: column; gap: 0.6rem; }}
    .toc-list li {{ font-size: 1.05rem; color: {subtitle_color}; cursor: default; transition: color 0.2s; }}
    .toc-list li:hover {{ color: {text}; }}
    .stat-grid {{
      display: grid; grid-template-columns: repeat(4,1fr);
      gap: 1px; background: {border}; border-radius: 10px; overflow: hidden;
      margin-top: 2rem; max-width: 600px;
    }}
    .stat-grid .stat {{ background: {stat_bg}; padding: 1.25rem 1rem; text-align: center; }}
    .stat .val {{ font-size: 1.15rem; font-weight: 500; color: {text}; }}
    .stat .lbl {{ font-size: 0.7rem; color: #87867f; margin-top: 0.15rem; }}
    .cmp-grid {{ display: grid; grid-template-columns: 1fr 1fr; gap: 1px; background: {border}; border-radius: 12px; overflow: hidden; margin-top: 1.5rem; }}
    .cmp-card {{ background: {card_bg}; padding: 2rem; }}
    .cmp-badge {{ display: inline-block; font-size: 0.7rem; font-weight: 500; text-transform: uppercase; letter-spacing: 0.12px; padding: 0.25rem 0.7rem; border-radius: 100px; margin-bottom: 1rem; background: {accent}; color: #faf9f5; }}
    .cmp-card p {{ font-size: 0.9rem; color: {subtitle_color}; line-height: 1.65; }}
    .gallery-grid {{ display: grid; grid-template-columns: repeat(2,1fr); gap: 1px; background: {border}; border-radius: 12px; overflow: hidden; margin-top: 2rem; }}
    .gallery-item {{ background: {card_bg}; aspect-ratio: 16/10; display: flex; flex-direction: column; align-items: center; justify-content: center; padding: 1rem; }}
    .gallery-placeholder {{ width: 80%; height: 60%; border-radius: 8px; margin-bottom: 0.75rem; }}
    .gallery-item .cap {{ font-size: 0.8rem; color: {subtitle_color}; }}
    .quote-slide {{ text-align: center; justify-content: center; align-items: center; }}
    .quote-text {{
      font-family: Georgia, 'Times New Roman', Times, serif;
      font-size: clamp(1.4rem,3vw,2.2rem); font-weight: 500;
      line-height: 1.6; color: {accent}; max-width: 700px;
      white-space: pre-line;
    }}
    .quote-attrib {{ margin-top: 1.5rem; color: #87867f; font-size: 0.9rem; }}
    .end {{ text-align: center; justify-content: center; align-items: center; }}
    .nav-bar {{ position: fixed; bottom: 2rem; right: 2rem; display: flex; gap: 0.4rem; z-index: 100; }}
    .nav-dot {{ width: 8px; height: 8px; border-radius: 50%; background: {border}; border: none; cursor: pointer; transition: background 0.2s; }}
    .nav-dot:hover {{ background: #87867f; }}
    .toolbar {{ position: fixed; top: 1rem; right: 1rem; z-index: 100; display: flex; gap: 0.5rem; }}
    .toolbar a {{
      display: inline-block; background: {card_bg}; border: 1px solid {border};
      color: {subtitle_color}; padding: 0.4rem 0.8rem; border-radius: 6px;
      text-decoration: none; font-size: 0.75rem; transition: all 0.2s;
    }}
    .toolbar a:hover {{ background: {border}; color: {text}; }}
    @media (max-width:640px) {{
      .slide {{ padding: 2rem 1.5rem; }}
      .stat-grid {{ grid-template-columns: repeat(2,1fr); }}
      .cmp-grid, .gallery-grid {{ grid-template-columns: 1fr; }}
    }}
  </style>
</head>
<body>
<div class="toolbar">
  <a href="#" onclick="window.print()">打印 / 导出 PDF</a>
  <a href="#" onclick="document.querySelector('.slide').scrollIntoView({{behavior:'smooth'}})">回到顶部</a>
</div>
<div class="nav-bar" id="navDots"></div>
{''.join(pages_html)}
<script>
  const slides = document.querySelectorAll('.slide');
  const nav = document.getElementById('navDots');
  slides.forEach((_,i) => {{
    const d = document.createElement('button');
    d.className='nav-dot';
    d.addEventListener('click',()=>slides[i].scrollIntoView({{behavior:'smooth'}}));
    nav.appendChild(d);
  }});
  const dots = nav.querySelectorAll('.nav-dot');
  new IntersectionObserver(entries => entries.forEach(e => {{
    if(e.isIntersecting) {{
      const idx = Array.from(slides).indexOf(e.target);
      dots.forEach((d,i)=>d.style.background=i===idx?'{accent}':'{border}');
    }}
  }}),{{threshold:0.6}}).observe(slides[0]);
  slides.forEach(s => new IntersectionObserver(entries => entries.forEach(e => {{
    if(e.isIntersecting) {{
      const idx = Array.from(slides).indexOf(e.target);
      dots.forEach((d,i)=>d.style.background=i===idx?'{accent}':'{border}');
    }}
  }}),{{threshold:0.6}}).observe(s));
</script>
</body>
</html>"""


def _render_page(page, idx, theme="dark"):
    t = page["type"]

    if t == "cover":
        return f"""<div class="slide cover">
  <div class="slide-label">建筑汇报方案</div>
  <h1>{page['title']}</h1>
  <div class="accent-line"></div>
  <p style="font-size:1.15rem;color:#87867f;">{page['subtitle']}</p>
  <div class="meta">{page['meta']}</div>
</div>"""

    if t == "toc":
        items = "".join(f"<li>{item}</li>" for item in page["items"])
        return f"""<div class="slide">
  <div class="slide-label">汇报框架</div>
  <h2>{page['title']}</h2>
  <ul class="toc-list" style="margin-top:2rem;">{items}</ul>
</div>"""

    if t == "content":
        stats_html = ""
        if page.get("stats"):
            stats = "".join(f'<div class="stat"><div class="val">{s["value"]}</div><div class="lbl">{s["label"]}</div></div>' for s in page["stats"])
            stats_html = f'<div class="stat-grid">{stats}</div>'
        return f"""<div class="slide">
  <div class="slide-label">第 {idx+1} 页</div>
  <h2>{page['title']}</h2>
  <p class="subtitle">{page['subtitle']}</p>
  <p>{page['body']}</p>
  {stats_html}
</div>"""

    if t == "split":
        return f"""<div class="slide">
  <div class="slide-label">第 {idx+1} 页</div>
  <h2>{page['title']}</h2>
  <div class="cmp-grid">
    <div class="cmp-card"><h3>{page.get('left_title','策略一')}</h3><p>{page['left']}</p></div>
    <div class="cmp-card"><h3>{page.get('right_title','策略二')}</h3><p>{page['right']}</p></div>
  </div>
</div>"""

    if t == "comparison":
        return f"""<div class="slide">
  <div class="slide-label">第 {idx+1} 页 · 方案对比</div>
  <h2>{page['title']}</h2>
  <p class="subtitle">{page['subtitle']}</p>
  <div class="cmp-grid">
    <div class="cmp-card">
      <div class="cmp-badge">{page.get('left_label','方案 A')}</div>
      <p>{page.get('left_body','')}</p>
    </div>
    <div class="cmp-card">
      <div class="cmp-badge">{page.get('right_label','方案 B')}</div>
      <p>{page.get('right_body','')}</p>
    </div>
  </div>
</div>"""

    if t == "gallery":
        items = ""
        for img in page.get("images", []):
            color = img.get("color", "#30302e")
            items += f"""<div class="gallery-item">
  <div class="gallery-placeholder" style="background:{color};"></div>
  <div class="cap">{img.get('caption','')}</div>
</div>"""
        return f"""<div class="slide">
  <div class="slide-label">第 {idx+1} 页 · 效果展示</div>
  <h2>{page['title']}</h2>
  <p class="subtitle">{page['subtitle']}</p>
  <div class="gallery-grid">{items}</div>
</div>"""

    if t == "quote":
        return f"""<div class="slide quote-slide">
  <div class="quote-text">{page['quote']}</div>
  <div class="quote-attrib">— {page.get('attribution','')}</div>
</div>"""

    if t == "end":
        return f"""<div class="slide end">
  <div class="accent-line" style="margin-bottom:1.5rem;"></div>
  <h2>{page['title']}</h2>
  <p class="subtitle" style="margin:0 auto;">{page['subtitle']}</p>
</div>"""

    return ""


# =============================================================================
# PPTX 导出 (实验性)
# =============================================================================

def try_export_pptx(html_path, pptx_path):
    """使用系统可用的工具转换为 PPTX。"""
    # 方法1: python-pptx (需先 pip install)
    try:
        import pptx
        return _export_pptx_python(html_path, pptx_path)
    except ImportError:
        pass

    # 方法2: pandoc
    try:
        subprocess.run(["pandoc", html_path, "-o", pptx_path], capture_output=True, timeout=30)
        if os.path.exists(pptx_path) and os.path.getsize(pptx_path) > 1000:
            return True
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    return False


def _export_pptx_python(html_path, pptx_path):
    """使用 python-pptx 生成基础 PPTX。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 解析 HTML 提取标题
    import re
    with open(html_path, "r", encoding="utf-8") as f:
        html = f.read()

    slides_parsed = re.findall(r'<h[12]>(.*?)</h[12]>', html)
    if not slides_parsed:
        slides_parsed = [os.path.splitext(os.path.basename(html_path))[0]]

    for title in slides_parsed:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x14, 0x14, 0x13)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title.strip()
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(0xFA, 0xF9, 0xF5)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Georgia"

    prs.save(pptx_path)
    return True


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(description="建筑汇报 HTML 预览生成器 v2.0")
    parser.add_argument("--out", "-o", default="preview.html", help="输出 HTML 路径")
    parser.add_argument("--json", "-j", help="从 JSON 文件读取页面数据")
    parser.add_argument("--theme", choices=["dark", "light"], default="dark", help="主题 (默认 dark)")
    parser.add_argument("--pptx", help="可选：导出 PPTX 文件路径")
    parser.add_argument("--list-types", action="store_true", help="列出所有支持的页面类型")
    args = parser.parse_args()

    if args.list_types:
        print("支持的页面类型：")
        for key, desc in PAGE_TYPES.items():
            print(f"  {key:15s} {desc}")
        return

    if args.json:
        with open(args.json, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = dict(DEFAULT_SECTIONS)

    data["theme"] = args.theme

    html = generate_html(data)
    out_path = os.path.abspath(args.out)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(html)

    print(f"✅ HTML 预览已生成: {out_path}")
    print(f"   共 {len(data['pages'])} 页 · 主题: {args.theme} · 项目: {data['title']}")

    if args.pptx:
        pptx_path = os.path.abspath(args.pptx)
        print(f"   尝试导出 PPTX...", end=" ")
        ok = try_export_pptx(out_path, pptx_path)
        if ok:
            print(f"✅ {pptx_path}")
        else:
            print("⚠️ 跳过（需要 python-pptx: pip install python-pptx）")


if __name__ == "__main__":
    main()
